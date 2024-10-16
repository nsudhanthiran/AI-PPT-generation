from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import ollama
import os

class Ppt:

    def getppt(prompt, no_slide):
        p1 = Presentation(r'D:\Code\Tot_PPT\static\ppt_temp\blueppt.pptx')
        desired_model = 'llama2:latest'
        
        # Generate titles using LLaMA model
        def title_for_content(prompt):
            response = ollama.chat(model=desired_model, messages=[
                {
                    'role': 'user',
                    'content': prompt,
                },
            ])
            return response['message']['content']

        # Generate content for each title using LLaMA model
        def content_on_title(prompt):
            response_content = ollama.chat(model=desired_model, messages=[
                {
                    'role': 'user',
                    'content': f'''Write content on the topic "{prompt}" in exactly 100 words. Do not include phrases like "Sure, here is the content" or mention the word count or any introductory lines. Start directly with the content.'''
                }
            ])
            return response_content['message']['content']

        # Adjust text alignment and font size
        def set_text_alignment_and_font(shape, alignment, font_size):
            for paragraph in shape.text_frame.paragraphs:
                paragraph.alignment = alignment  # Align text (left, center, right)
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)  # Set the font size dynamically
        
        # Alter headings
        def alter_headings(heading=filter_heading):
            head = []
            for i in heading:
                head.append(i)
            return head


        # Create PowerPoint slides
        def create_ppt(titles_raw):
            for title in titles_raw:
                content = content_on_title(title)
                
                # Use the slide layout from the template
                slide_register = p1.slide_layouts[1]
                slide = p1.slides.add_slide(slide_register)
                
                # Access title and content placeholders
                slide_title = slide.shapes.title
                slide_content = slide.placeholders[1]

                # Insert title and content text
                slide_title.text = title
                slide_content.text = content

                # Set alignment and font size for title and content
                set_text_alignment_and_font(slide_title, PP_ALIGN.CENTER, 32)  # Center align title with larger font
                set_text_alignment_and_font(slide_content, PP_ALIGN.LEFT, 20)  # Left align content with smaller font

            # Save the newly created PowerPoint
            direction = r"D:\Code\Tot_PPT\static\ppt_temp"
            path_dir = os.path.join(direction, 'new.pptx')
            p1.save(path_dir)

        # Generate the titles for the slides
        headings = title_for_content(prompt=f'''Give me {no_slide} titles for the keyword '{prompt}' without any additional text or explanations. Respond with just the titles.''')
        
        # Process the headings to filter out empty lines
        nor_heading = f'''{headings}'''.splitlines()
        filter_heading = [title for title in nor_heading if title.strip()]

        alter_headings(filter_heading)

        # Create the PowerPoint based on the generated titles and content
        create_ppt(titles_raw=filter_heading)

